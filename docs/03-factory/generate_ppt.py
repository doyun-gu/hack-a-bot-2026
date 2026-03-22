"""
Generate GridBox hackathon presentation PPT.
Usage: python3 docs/03-factory/generate_ppt.py
Output: docs/03-factory/GridBox-Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colours
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)
MID_BLUE = RGBColor(0x16, 0x21, 0x3e)
ACCENT_RED = RGBColor(0xe9, 0x45, 0x60)
ACCENT_GREEN = RGBColor(0x2d, 0x6a, 0x4f)
ACCENT_PURPLE = RGBColor(0x53, 0x34, 0x83)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK_BLUE):
    """Fill slide background with solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, font_name="Calibri"):
    """Add bulleted list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return tf

def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=ACCENT_PURPLE, cell_color=MID_BLUE):
    """Add a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = WHITE
                paragraph.font.name = "Calibri"

            # Header row styling
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_color

    return table

def add_shape_box(slide, left, top, width, height, text, fill_color=ACCENT_PURPLE,
                  font_size=14):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_text_box(slide, 1, 1.5, 11, 1.5, "GridBox", font_size=72, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.0, 11, 1, "Smart Infrastructure Control for £15",
             font_size=32, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 11, 0.8,
             '"What if a £15 microcontroller could replace a £162,000 industrial system?"',
             font_size=20, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.5, 11, 0.5, "ARM + EEESoc  |  Hack-A-Bot 2026",
             font_size=18, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.1, 11, 0.5, "Doyun Gu  ·  Wooseong  ·  Billy",
             font_size=16, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

# Try to add UoM logo
logo_path = os.path.expanduser("~/Documents/02-EEE/UoM-logo/png/TAB_allwhite.png")
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(10.5), Inches(0.3), Inches(2.2))


# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, 0.8, 0.5, 11, 1, "The Problem", font_size=44, bold=True, color=ACCENT_RED)

# Three problem boxes
add_shape_box(slide, 0.8, 2.0, 3.6, 2.5,
              "30%\nof factory energy\nis wasted\n\nMotors run at full speed\nwhen partial load is enough",
              fill_color=ACCENT_RED, font_size=16)

add_shape_box(slide, 4.8, 2.0, 3.6, 2.5,
              "£180B/year\nlost to undetected\nequipment faults\n\n4+ hours average\ndowntime per fault",
              fill_color=ACCENT_PURPLE, font_size=16)

add_shape_box(slide, 8.8, 2.0, 3.6, 2.5,
              "£162,000+\ncost of industrial\nSCADA system\n\n80% of small factories\npriced out",
              fill_color=MID_BLUE, font_size=16)

add_text_box(slide, 0.8, 5.2, 11, 1.5,
             "Small and medium factories can't afford Industry 4.0.\n"
             "They run blind — no monitoring, no fault detection, no optimisation.\n"
             "When equipment fails, they find out the hard way.",
             font_size=18, color=LIGHT_GREY)


# ============================================================
# SLIDE 3: Our Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, 0.8, 0.5, 11, 1, "Our Solution — GridBox", font_size=44, bold=True, color=ACCENT_GREEN)

# Architecture boxes
add_shape_box(slide, 0.8, 1.8, 3.8, 1.5,
              "PICO A — Master\nGrid Controller\n\nIMU · PCA9685 · ADC\nMotor driver · Recycle path",
              fill_color=ACCENT_PURPLE, font_size=14)

add_shape_box(slide, 5.0, 2.0, 2.8, 1.1,
              "nRF24L01+ Wireless\n2.4GHz · 250kbps\n32-byte datagrams",
              fill_color=ACCENT_RED, font_size=13)

add_shape_box(slide, 8.5, 1.8, 3.8, 1.5,
              "PICO B — Slave\nSCADA Station\n\nOLED · MAX7219 · nRF\nWireless dashboard",
              fill_color=MID_BLUE, font_size=14)

# Closed loop
add_text_box(slide, 0.8, 3.8, 11, 0.8, "Closed-Loop Control (every 10ms):",
             font_size=20, bold=True, color=WHITE)
add_text_box(slide, 0.8, 4.5, 11, 0.8,
             "SENSE (ADC)  →  CALCULATE (KCL/KVL)  →  DECIDE (firmware)  →  ROUTE (GPIO)  →  VERIFY (ADC)",
             font_size=18, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Demo scenario
add_text_box(slide, 0.8, 5.5, 11, 0.5, "Demo: Smart Water Bottling Plant",
             font_size=20, bold=True, color=WHITE)

demo_data = [
    ["Stage", "Hardware", "What Happens"],
    ["Pump", "DC Motor 1", "Fills bottles at PWM-controlled speed"],
    ["Fill Valve", "Servo 1 (PCA9685)", "Opens/closes water flow"],
    ["Conveyor", "DC Motor 2", "Moves bottles along the line"],
    ["Sort Gate", "Servo 2 (PCA9685)", "Directs good/reject bottles"],
    ["Fault Detect", "BMI160 IMU", "Vibration monitoring (ISO 10816)"],
    ["Energy Recycle", "2N2222 + Cap + LED", "Captures and reuses wasted energy"],
]
add_table(slide, 0.8, 6.0, 11.5, 1.2, len(demo_data), 3, demo_data)


# ============================================================
# SLIDE 4: Technical Innovation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, 0.8, 0.5, 11, 1, "Technical Innovation", font_size=44, bold=True, color=ORANGE)

# Core innovation
add_shape_box(slide, 0.8, 1.7, 5.5, 1.8,
              "The Pico IS the Power Grid\n\n"
              "ADC pins = power sensors\n"
              "GPIO pins = electronic switches\n"
              "Firmware = brain that decides\n"
              "Closed loop = autonomous control",
              fill_color=ACCENT_PURPLE, font_size=15)

# EEE Theory
add_shape_box(slide, 6.8, 1.7, 5.5, 1.8,
              "EEE Theory Applied\n\n"
              "Affinity Laws: P ∝ n³\n"
              "KCL/KVL energy balance\n"
              "ISO 10816 vibration classification\n"
              "Current sensing: I = V/R",
              fill_color=MID_BLUE, font_size=15)

# Three innovation pillars
add_shape_box(slide, 0.8, 4.0, 3.6, 2.2,
              "IMU Fault Detection\n\n"
              "£2 BMI160 replaces\n£18,000 equipment\n\n"
              "Detects: bearing wear,\ncavitation, jams, leaks",
              fill_color=ACCENT_RED, font_size=13)

add_shape_box(slide, 4.8, 4.0, 3.6, 2.2,
              "Energy Recycling\n\n"
              "GPIO routes excess\npower to storage\n\n"
              "LED proves concept:\ncharge → glow → fade",
              fill_color=ACCENT_GREEN, font_size=13)

add_shape_box(slide, 8.8, 4.0, 3.6, 2.2,
              "Wireless SCADA\n\n"
              "6-type datagram protocol\n200+ packets, 0 errors\n\n"
              "Bidirectional:\ncommands + telemetry",
              fill_color=ACCENT_PURPLE, font_size=13)

# Key equation
add_text_box(slide, 0.8, 6.5, 11, 0.7,
             "Key insight: 20% speed reduction = 49% power saving (Affinity Laws: P ∝ n³)",
             font_size=18, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: What We Built + Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, 0.8, 0.5, 11, 1, "What We Built", font_size=44, bold=True, color=ACCENT_GREEN)

# Results table
results_data = [
    ["Metric", "Result", "Significance"],
    ["Cost", "£15 total BOM", "10,800x cheaper than industrial SCADA"],
    ["Wireless", "200+ packets, 0 errors", "6 datagram types, bidirectional"],
    ["Wiring", "77% complete (48/66 wires)", "10/13 tasks done"],
    ["Firmware", "28 MicroPython modules", "16 master + 12 slave + C SDK"],
    ["Tests", "35+ test scripts", "Every component verified individually"],
    ["Docs", "35+ documents", "Full design doc, wiring tables, diagrams"],
    ["Commits", "210+ commits", "3-person team, 24-hour hackathon"],
]
add_table(slide, 0.8, 1.7, 11.5, 2.5, len(results_data), 3, results_data,
          header_color=ACCENT_GREEN)

# Hardware verified
add_text_box(slide, 0.8, 4.5, 11, 0.5, "Hardware Verified:",
             font_size=20, bold=True, color=WHITE)

verified_items = [
    "nRF24L01+ wireless (both Picos) — SPI status 0x0E, channel read/write OK",
    "Two-Pico wireless link — standalone master + USB slave, datagram protocol",
    "MAX7219 7-segment display — all digits, brightness, flash animations",
    "BMI160 IMU — I2C connected, chip ID verified",
    "Motor driver — connected and wired via PCA9685 PWM",
    "Recycle path — 2N2222 + capacitor + LED discharge indicator",
]
add_bullet_list(slide, 0.8, 5.0, 11, 2.0, verified_items, font_size=14, color=LIGHT_GREY)

# Cost comparison at bottom
add_shape_box(slide, 3.5, 6.8, 6, 0.5,
              "£15  vs  £162,000  —  Same job, 10,800x cheaper",
              fill_color=ACCENT_RED, font_size=16)


# ============================================================
# SLIDE 6: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, 1, 2.0, 11, 1.5, "GridBox", font_size=72, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.5, 11, 1,
             "Smart Infrastructure Control for £15",
             font_size=28, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.8, 11, 1,
             "Thank you — we're happy to take questions.",
             font_size=22, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 6.0, 11, 0.5,
             "Doyun Gu  ·  Wooseong  ·  Billy  |  University of Manchester",
             font_size=16, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)

if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(5.5), Inches(6.5), Inches(2.2))


# ============================================================
# Save
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "GridBox-Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
